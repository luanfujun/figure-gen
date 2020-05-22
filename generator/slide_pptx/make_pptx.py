import os
import json
from pptx import Presentation
from pptx.util import Inches
from . import place_element, calculate

'''
in PPTX format we 
 - ignore background colors
 - ignore element captions (north/east/south/west content of each img) as we didn't even use them once before
 - do not support 'dashed' frames - if a frame is 'dashed' the frame in pptx will ne normal (but still has a frame)
''' 
    
def generate(module_data, to_path, index, delete_gen_files=True):
    return module_data

def combine(data, to_path, delete_gen_files=True):
    #create slide
    prs = Presentation()
    prs.slide_height = Inches(9) 
    prs.slide_width = Inches(16)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # calculate correct width scaling so that the figure fills out the slide
    sum_total_width_mm = 0
    for d in data:
        sum_total_width_mm += d['total_width']
    width_scaling = 16 / calculate.mm_to_inch(sum_total_width_mm)

    # generate content
    cur_width_mm = 0
    for d in data:
        place_element.images_and_frames(slide, d, width_scaling, cur_width_mm)
        cur_width_mm += d['total_width']

    # save
    path_file = os.path.join(to_path, 'gen_figure.pptx')
    prs.save(path_file)